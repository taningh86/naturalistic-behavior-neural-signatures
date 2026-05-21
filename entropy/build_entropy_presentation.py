"""
Build PowerPoint: 'Behavioral Entropy and its Neural Correlates'
15 slides max; large legible text; brief figure descriptions.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

REPO = Path(r"H:\NPX ANALYSIS REPO")
FIG = REPO / "figures"
OUT = REPO / "data" / "behavior_entropy_and_neural_correlates.pptx"

# Colors
NAVY = RGBColor(0x0B, 0x2E, 0x4F)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    return s


def add_title_bar(slide, text, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(16)
        r2.font.color.rgb = RGBColor(0xD5, 0xE0, 0xF0)


def add_text(slide, left, top, width, height, lines, size=18,
             bold_first=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    if not isinstance(lines, (list, tuple)):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        if bold_first and i == 0:
            run.font.bold = True
        run.font.color.rgb = color if color else RGBColor(0x20, 0x20, 0x20)
    return tb


def add_caption(slide, left, top, width, height, text, size=16):
    """Figure caption box — italicized, gray background."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(0x20, 0x20, 0x20)


def add_figure(slide, path, left, top, width, height):
    p = Path(path)
    if not p.exists():
        # placeholder box if missing
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(0xFA, 0xE5, 0xE5)
        b.text_frame.text = f"[missing: {p.name}]"
        return None
    return slide.shapes.add_picture(str(p), left, top, width=width, height=height)


def add_bullet_block(slide, left, top, width, height, bullets, size=18):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


# ============================================================
# SLIDE 1 — Title
# ============================================================
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "Behavioral Entropy and its Neural Correlates"
r.font.size = Pt(46)
r.font.bold = True
r.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Shannon entropy of zone transitions as a behavioral readout, and how LHA, RSP, and ACA track its peaks and troughs"
r.font.size = Pt(22)
r.font.color.rgb = RGBColor(0xD5, 0xE0, 0xF0)

tb = s.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Neuropixels foraging project  |  single- and dual-probe summary"
r.font.size = Pt(16)
r.font.color.rgb = RGBColor(0xB8, 0xC4, 0xD6)

# ============================================================
# SLIDE 2 — Shannon entropy
# ============================================================
s = add_slide()
add_title_bar(s, "Shannon Entropy: Measuring Unpredictability")

# Formula box
fb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                       Inches(0.8), Inches(1.2),
                       Inches(5.8), Inches(2.4))
fb.fill.solid()
fb.fill.fore_color.rgb = LIGHT
fb.line.color.rgb = NAVY
tf = fb.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "H(X) = - Σ  p(x) log₂ p(x)"
r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = NAVY
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "Units: bits"
r.font.size = Pt(20); r.font.italic = True; r.font.color.rgb = GRAY

add_bullet_block(s, Inches(7.1), Inches(1.2), Inches(5.8), Inches(4.0), [
    "H = 0 bits → outcome is fully predictable",
    "H is maximal when all outcomes are equally likely",
    "High H → disordered, varied, exploratory",
    "Low H  → repetitive, stereotyped, habitual",
    "We compute H on the distribution of zone-to-zone transitions",
], size=20)

add_caption(s, Inches(0.8), Inches(5.6), Inches(12.0), Inches(1.3),
    "Intuition: a fair coin has H = 1 bit; a mouse that moves Home→Ladder→Home→Ladder has nearly zero entropy; "
    "a mouse visiting many zones in varied order has high entropy.")

# ============================================================
# SLIDE 3 — Entropy applied to mouse zone transitions
# ============================================================
s = add_slide()
add_title_bar(s, "From Zone Visits to an Entropy Trace")

add_bullet_block(s, Inches(0.6), Inches(1.1), Inches(12.2), Inches(3.6), [
    "Arena divided into zones: Home, Ladder, T-zone, Pot-1..4, seed pot, transition wall",
    "For every 60-s sliding window (10-s step), we count each zone → zone transition",
    "Shannon entropy of this transition distribution = behavioral entropy at that moment",
    "Output: one entropy value per 10 s throughout the 30-min session",
    "Complementary to zone occupancy: captures WHICH pair of zones was traversed, not just where the mouse was",
], size=20)

add_caption(s, Inches(0.6), Inches(5.4), Inches(12.2), Inches(1.7),
    "Why transitions (not zone occupancy)? A mouse that sits in one zone for a long time can still show structured transitions "
    "nearby. Transition entropy specifically measures the predictability of the animal's next movement choice.",
    size=17)

# ============================================================
# SLIDE 4 — Entropy across conditions
# ============================================================
s = add_slide()
add_title_bar(s, "Entropy Across Metabolic State x Task Phase")

add_figure(s, FIG / "behav_entropy_all_conditions.png",
           Inches(0.4), Inches(1.05), Inches(8.3), Inches(5.6))

add_bullet_block(s, Inches(8.9), Inches(1.2), Inches(4.2), Inches(5.0), [
    "Fed / Fasted / HFD x Exp / Foraging",
    "Fasted + foraging has LOWEST mean entropy and deepest dips (to ~1 bit)",
    "Fasted + exploration stays high — fasting alone doesn't reduce entropy",
    "Task demand + hunger together drive stereotyped foraging",
], size=17)

add_caption(s, Inches(0.4), Inches(6.75), Inches(12.6), Inches(0.6),
    "Behavioral entropy traces pooled across all 20 analyzed dual-probe sessions, split by metabolic state and task phase.",
    size=15)

# ============================================================
# SLIDE 5 — Broader behavioral differences Fed vs Fasted vs HFD
# ============================================================
s = add_slide()
add_title_bar(s, "Behavioral Differences: Fed vs Fasted vs HFD",
              subtitle="20 metrics across 20 dual-probe sessions")

add_figure(s, FIG / "dp_behavior_hfd_comparison.png",
           Inches(0.25), Inches(1.0), Inches(9.3), Inches(6.2))

add_bullet_block(s, Inches(9.7), Inches(1.15), Inches(3.5), Inches(5.8), [
    "FASTED:",
    "  • LESS digging (fewer bouts, shorter)",
    "  • LOWER velocity",
    "  • MORE immobility",
    "",
    "HFD — UNIQUE signature:",
    "  • ↑ transition-zone occupancy",
    "  • ↑ total pot visits (210 vs 168/105)",
    "  • ↑ incomplete returns",
    "  • ↑ total distance",
    "",
    "HFD pattern:",
    "  restless / indecisive foraging",
], size=14)

add_caption(s, Inches(0.25), Inches(7.0), Inches(12.8), Inches(0.4),
    "Bars: mean ± SEM across sessions; dots: individual sessions. KW p-values shown on significant panels.",
    size=12)

# ============================================================
# SLIDE 6 — Example session with behavioral annotations
# ============================================================
s = add_slide()
add_title_bar(s, "Example Session: Entropy Aligned with Scored Behavior",
              subtitle="S6 — fed / foraging")

add_figure(s, FIG / "entropy_behavior_alignment" / "S6_fed_foraging_entropy_behavior.png",
           Inches(0.3), Inches(1.05), Inches(9.3), Inches(5.6))

add_bullet_block(s, Inches(9.8), Inches(1.2), Inches(3.4), Inches(5.0), [
    "Blue: raw entropy",
    "Black: smoothed",
    "Orange (right axis): velocity",
    "Below: manually-scored behaviors",
    "Troughs coincide with feeding / digging bouts",
    "Peaks coincide with transitions between zones",
], size=17)

add_caption(s, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.6),
    "Each trough parks the mouse in one compartment and one behavior (dig / feed). Peaks are the travel between them.",
    size=15)

# ============================================================
# SLIDE 6 — Behaviors driving peaks vs troughs
# ============================================================
s = add_slide()
add_title_bar(s, "What Drives Peaks vs Troughs?",
              subtitle="Behavioral modes at low entropy across all 8 single-probe sessions")

add_figure(s, FIG / "entropy_mode_neural_traces.png",
           Inches(0.3), Inches(1.05), Inches(9.3), Inches(5.6))

add_bullet_block(s, Inches(9.8), Inches(1.2), Inches(3.4), Inches(5.2), [
    "PEAKS → varied zone traversal, exploration",
    "TROUGHS → one of three motifs:",
    "  • H/L shuttling (fed)",
    "  • P2 loops (exploration)",
    "  • P4 loops (fasted foraging)",
    "Distinct stereotyped modes, not a single one",
], size=16)

add_caption(s, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.6),
    "Low-entropy windows were categorized by dominant zone-pair transitions (>35% share). Different states use different motifs.",
    size=15)

# ============================================================
# SLIDE 7 — Single-probe LHA-RSP correlation heatmap
# ============================================================
s = add_slide()
add_title_bar(s, "Single Probe: LHA and RSP Show Opposing Entropy Correlations")

add_figure(s, FIG / "entropy_neural_correlation_summary.png",
           Inches(0.3), Inches(1.05), Inches(8.5), Inches(5.6))

add_bullet_block(s, Inches(9.0), Inches(1.2), Inches(4.2), Inches(5.2), [
    "RSP FR / PC1: POSITIVE with entropy (7/8 sessions)",
    "LHA FR / PC1: NEGATIVE with entropy (6/8 sessions)",
    "43/64 Spearman tests significant (raw)",
    "Strongest: S6 (fasted foraging)",
    "Interpretation: LHA ↑ during stereotypy, RSP ↑ during varied behavior",
], size=17)

add_caption(s, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.6),
    "Each cell = Spearman rho between a neural metric and behavioral entropy. Warm = positive, cool = negative.",
    size=15)

# ============================================================
# SLIDE 8 — Single-probe ±120 s pooled peri-inflection
# ============================================================
s = add_slide()
add_title_bar(s, "Single Probe: ±120 s Peri-Inflection — LHA ↔ RSP Opposition",
              subtitle="Pooled across 8 sessions; z-scored to −120 to −60 s baseline")

add_figure(s, FIG / "entropy_inflection_pooled_key.png",
           Inches(0.3), Inches(1.05), Inches(8.5), Inches(5.4))

add_bullet_block(s, Inches(9.0), Inches(1.15), Inches(4.2), Inches(5.6), [
    "Pooled peaks and troughs across 8 sessions",
    "RSP FR/PC1 RISES toward peaks, FALLS toward troughs",
    "LHA FR/PC1 FALLS toward peaks, RISES toward troughs",
    "Within-session validation:",
    "  • RSP FR: 8/8 sessions correct direction",
    "     Wilcoxon p=0.008**",
    "  • RSP PC1: 8/8 sessions; p=0.008**",
    "  • LHA PC1: 7/8 sessions; p=0.11 (borderline)",
    "Effect lives on a 60–120 s timescale",
], size=14)

add_caption(s, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.75),
    "Mean ± SEM across pooled events. Neural swings build up over tens of seconds around each entropy inflection — "
    "a slow covariation, not a transient event lock.",
    size=14)

# ============================================================
# SLIDE 9 — Single-probe tight peri-inflection (honest null)
# ============================================================
s = add_slide()
add_title_bar(s, "Single Probe: Tight ±5 s Peri-Inflection — No Event-Locked Effect",
              subtitle="500 ms bins, z-scored to −10 s to −5 s baseline")

add_figure(s, FIG / "entropy_inflection_tight_sp.png",
           Inches(0.3), Inches(1.15), Inches(8.5), Inches(4.8))

add_bullet_block(s, Inches(9.0), Inches(1.2), Inches(4.2), Inches(5.4), [
    "Pooled: 34 peaks, 38 troughs (8 sessions)",
    "At PEAKS:",
    "  • Velocity p=0.11 ns",
    "  • LHA FR p=0.12 ns",
    "  • RSP FR p=0.45 ns",
    "  • LHA PC1 p=0.47 ns",
    "  • RSP PC1 p=0.24 ns",
    "No MWU peak-vs-trough significant",
    "Traces drift monotonically — no tight event lock",
], size=15)

add_caption(s, Inches(0.3), Inches(6.05), Inches(12.7), Inches(1.2),
    "When restricted to ±5 s, the single-probe LHA-RSP opposition does NOT reach significance. "
    "The apparent 120-s opposition is a slow covariation, not an event-locked transient. Direction still weakly "
    "consistent (LHA FR trends up at peaks) but statistical power is exhausted.",
    size=14)

# ============================================================
# SLIDE 9 — Dual probe: setup + correlations
# ============================================================
s = add_slide()
add_title_bar(s, "Dual Probe: ACA Takes RSP's Role Opposite LHA")

add_figure(s, FIG / "dp_entropy_neural_correlations.png",
           Inches(0.3), Inches(1.05), Inches(8.5), Inches(5.6))

add_bullet_block(s, Inches(9.0), Inches(1.2), Inches(4.2), Inches(5.2), [
    "Probe 0 = ACA; Probe 1 = LHA",
    "18 sessions: fed S3–S10, fasted S11–S16, HFD S19–S22",
    "ACA FR: positive with entropy",
    "LHA FR/PC1: negative with entropy",
    "Fasted state: STRONGEST opposition",
    "Velocity again attenuates raw correlations",
], size=17)

add_caption(s, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.6),
    "Steady-state Spearman rho per session. Opposition is clearest in fasted sessions; HFD shows a direction flip.",
    size=15)

# ============================================================
# SLIDE 11 — Dual-probe ±120 s pooled peri-inflection
# ============================================================
s = add_slide()
add_title_bar(s, "Dual Probe: ±120 s Peri-Inflection — ACA ↔ LHA Opposition",
              subtitle="76 peaks, 81 troughs pooled across 17 sessions")

add_figure(s, FIG / "dp_entropy_inflection_pooled.png",
           Inches(0.3), Inches(1.05), Inches(8.5), Inches(5.4))

add_bullet_block(s, Inches(9.0), Inches(1.15), Inches(4.2), Inches(5.6), [
    "ACA FR RISES toward peaks, FALLS toward troughs",
    "LHA FR FALLS toward peaks, RISES toward troughs",
    "LHA leads troughs by ~30–40 s",
    "Within-session validation:",
    "  • ACA FR: 12/17 (71%) direction correct",
    "     Wilcoxon p=0.035*",
    "  • LHA FR: 10/17 (59%) ns",
    "Opposition mirrors the single-probe LHA-RSP result",
    "Operates on a slow (seconds to tens of seconds) timescale",
], size=14)

add_caption(s, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.75),
    "Entropy, velocity, ACA and LHA FR + PC1 z-scored to pre-inflection baseline. Neural swings are confined to "
    "~±60 s of the inflection.",
    size=14)

# ============================================================
# SLIDE 12 — Dual-probe tight peri-inflection
# ============================================================
s = add_slide()
add_title_bar(s, "Dual Probe: Tight ±5 s Peri-Inflection — Only Velocity Event-Locked",
              subtitle="90 peaks, 89 troughs pooled across 18 sessions; baseline −10 s to −5 s")

add_figure(s, FIG / "dp_entropy_inflection_tight.png",
           Inches(0.3), Inches(1.15), Inches(8.5), Inches(4.8))

# Stats table
tb = s.shapes.add_textbox(Inches(9.0), Inches(1.2), Inches(4.2), Inches(5.4))
tf = tb.text_frame
tf.word_wrap = True
lines = [
    ("At PEAKS (entropy high)", True),
    ("  • Velocity p=0.0004***", False),
    ("  • ACA FR p=0.89 ns", False),
    ("  • LHA FR p=0.16 ns", False),
    ("  • ACA PC1 p=0.20 ns", False),
    ("  • LHA PC1 p=0.29 ns", False),
    ("", False),
    ("At TROUGHS (entropy low)", True),
    ("  • Velocity p=0.09 ns", False),
    ("  • ACA FR p=0.27 ns", False),
    ("  • LHA FR p=0.37 ns", False),
    ("  • ACA PC1 p=0.69 ns", False),
    ("  • LHA PC1 p=0.49 ns", False),
]
for i, (text, bold) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.size = Pt(15) if bold else Pt(13)
    r.font.bold = bold
    r.font.color.rgb = NAVY if bold else RGBColor(0x22, 0x22, 0x22)

add_caption(s, Inches(0.3), Inches(6.05), Inches(12.7), Inches(1.2),
    "Only velocity shows a tight event lock (subtle dip around peaks). Pooled firing-rate and PC1 changes fail "
    "significance at ±5 s — the ACA-LHA opposition operates on a slower timescale (seconds to tens of seconds) "
    "rather than at the inflection instant.",
    size=14)

# ============================================================
# SLIDE 13 — Dual-probe ±120 s by state
# ============================================================
s = add_slide()
add_title_bar(s, "Dual Probe ±120 s by State: Fasted Strongest, HFD Abolishes Opposition",
              subtitle="Fed / Fasted / HFD peri-inflection at 120-s window")

add_figure(s, FIG / "dp_entropy_inflection_by_state.png",
           Inches(0.3), Inches(1.05), Inches(8.5), Inches(5.4))

add_bullet_block(s, Inches(9.0), Inches(1.15), Inches(4.2), Inches(5.6), [
    "FED: ACA/LHA opposition present",
    "FASTED: STRONGEST ACA/LHA opposition",
    "  • pooled ACA FR, LHA FR both sig",
    "HFD: opposition ABSENT / REVERSED",
    "  • 0–1 of 4 HFD sessions consistent",
    "Consistent with state-dependent reorganization",
    "HFD may uncouple the ACA↔LHA circuit",
], size=15)

add_caption(s, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.75),
    "Metabolic state gates the ACA-LHA opposition. Fasted sessions show the largest neural swings around entropy "
    "inflections; HFD shows flat or inverted patterns.",
    size=14)

# ============================================================
# SLIDE 14 — Tight window by state
# ============================================================
s = add_slide()
add_title_bar(s, "State Dependence at ±5 s: Scattered Effects, No Consistent Opposition",
              subtitle="fed / fasted / fed-HFD tight peri-inflection")

add_figure(s, FIG / "dp_entropy_inflection_tight_by_state.png",
           Inches(0.3), Inches(1.15), Inches(8.5), Inches(4.8))

add_bullet_block(s, Inches(9.0), Inches(1.2), Inches(4.2), Inches(5.4), [
    "FED (40 peaks / 40 troughs)",
    "  Velocity peak p=0.027*, ACA PC1 peak p=0.030*",
    "  All FR metrics ns",
    "",
    "FASTED (29 / 27)",
    "  LHA PC1 peak p=0.080 borderline (+6.48 z)",
    "  All others ns",
    "",
    "FED-HFD (21 / 22)",
    "  Velocity peak p=0.018*, MWU p=0.038*",
    "  All neural metrics ns",
], size=14)

add_caption(s, Inches(0.3), Inches(6.05), Inches(12.7), Inches(1.2),
    "At this tight timescale, only velocity reaches significance consistently (fed, HFD peaks). The 120-s opposition "
    "between ACA and LHA does not replicate as a transient ±5-s effect. Pooled LHA PC1 rise in fasted trends in the "
    "expected direction (+6 z) but fails Wilcoxon.",
    size=14)

# ============================================================
# SLIDE 12 — Phase effect (Exp vs Foraging)
# ============================================================
s = add_slide()
add_title_bar(s, "Phase Effect: Foraging Pulls Entropy Lower, Opposition Stronger")

add_figure(s, FIG / "behav_entropy_fasted_exp_vs_for.png",
           Inches(0.3), Inches(1.05), Inches(6.6), Inches(5.6))

add_bullet_block(s, Inches(7.1), Inches(1.15), Inches(6.0), Inches(5.6), [
    "Exploration (odd sessions):",
    "  • food visible in seed pot-2",
    "  • entropy stays HIGH (~4.2–4.6 bits)",
    "",
    "Foraging (even sessions):",
    "  • food hidden in sand pot-3 or pot-4",
    "  • entropy DIPS toward ~1 bit",
    "",
    "Neural opposition amplified during foraging:",
    "  • LHA engages stable attractor",
    "  • ACA / RSP tracks the variable search",
], size=17)

add_caption(s, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.6),
    "Task phase alone (hidden vs visible food) is enough to shift entropy distribution; amplified when combined with fasting.",
    size=15)

# ============================================================
# SLIDE 13 — Temporal dynamics
# ============================================================
s = add_slide()
add_title_bar(s, "Entropy Dips Are Oscillations, Not Single State Switches")

add_figure(s, FIG / "entropy_temporal_dynamics.png",
           Inches(0.3), Inches(1.05), Inches(8.5), Inches(5.6))

add_bullet_block(s, Inches(9.0), Inches(1.2), Inches(4.2), Inches(5.2), [
    "83–100% of entropy variance is OSCILLATION",
    "S6 has ZERO linear trend (p=0.96)",
    "Entropy drops are GRADUAL (autocorr 0.80–0.86)",
    "Many peaks/troughs per session → provides statistical power for pooled peri-analysis",
], size=18)

add_caption(s, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.6),
    "Recurring peaks and troughs (not a one-shot shift) is what made the 76-peak / 81-trough pooled test possible.",
    size=15)

# ============================================================
# SLIDE 14 — Summary
# ============================================================
s = add_slide()
add_title_bar(s, "Summary: Behavioral Entropy Has a Neural Signature")

bullets = [
    "Shannon entropy of zone transitions tracks foraging predictability in real time",
    "Troughs are STATE-SPECIFIC stereotypy motifs (H/L, P2 loops, P4 loops); peaks are varied traversal",
    "Single probe: LHA ↑ during troughs, RSP ↑ during peaks on 120-s windows — movement-confounded in raw correlation",
    "Dual probe: ACA substitutes for RSP; ACA-LHA opposition holds on 120-s windows (p < 0.01 at peaks and troughs)",
    "Tight ±5 s peri-inflection: only velocity event-locked; firing-rate / PC1 opposition operates on slower timescales",
    "State dependence at slow scale: Fasted > Fed ≫ HFD (HFD abolishes the opposition)",
]
add_bullet_block(s, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.6), bullets, size=22)

add_caption(s, Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.7),
    "Take-home: LHA and cortical (RSP / ACA) populations encode orthogonal axes of behavioral structure — "
    "not just movement.",
    size=18)

# ============================================================
# SLIDE 15 — Caveats and next steps
# ============================================================
s = add_slide()
add_title_bar(s, "Caveats and Next Steps")

# two columns
cap_left = s.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(6.0), Inches(5.8))
tf = cap_left.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Caveats"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = ACCENT
for b in [
    "Single-probe raw correlations are largely velocity-driven",
    "Dual-probe per-session tests are underpowered (0–3/17 sig); pooling produces the 120-s effect",
    "Tight ±5 s peri-inflection is null for neural metrics — the opposition is a slow covariation, not a transient",
    "HFD null result could be biology OR limited N (4 sessions)",
    "Entropy operates on a 60-s sliding window — coarse timescale by construction",
    "Manual behavior labels missing for fed S5–S10 (partial coverage)",
]:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "•  " + b
    r.font.size = Pt(17); r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    p.space_after = Pt(6)

cap_right = s.shapes.add_textbox(Inches(7.0), Inches(1.1), Inches(6.0), Inches(5.8))
tf = cap_right.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Next Steps"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = NAVY
for b in [
    "GRU-ODE mechanistic metrics already show S6 fasted/foraging as the inflection-point session",
    "Test at matched velocity bins to strengthen causal claim",
    "Extend dual-probe analysis once more fasted sessions are collected",
    "Investigate HFD-specific reorganization — why does opposition vanish?",
    "Relate entropy inflections to specific transition events (incomplete returns, T-zone contemplation)",
]:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "•  " + b
    r.font.size = Pt(17); r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    p.space_after = Pt(6)


# Save
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
