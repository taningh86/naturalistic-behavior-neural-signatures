"""
Build PowerPoint presentation: Foraging Neural Signatures.

Sections:
1. Pot-2 pre and post-discovery
2. Pot-4 pre and post-discovery
3. P2<->P4 transitions
+ Synthesis and take-home
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FIG = Path("figures")
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
GOLD = RGBColor(0xF3, 0x9C, 0x12)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=12,
                   color=BLACK, font_name='Calibri', spacing=1.0):
    """lines is list of (text, bold, color_override) tuples or just strings."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, bold, col = line, False, color
        else:
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            col = line[2] if len(line) > 2 else color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = col
        p.font.name = font_name
        p.space_after = Pt(font_size * spacing * 0.5)
    return txBox


def add_image(slide, img_path, left, top, width=None, height=None):
    path = str(img_path)
    if not os.path.exists(path):
        print(f"  WARNING: {path} not found, skipping")
        return None
    kwargs = {'left': Inches(left), 'top': Inches(top)}
    if width:
        kwargs['width'] = Inches(width)
    if height:
        kwargs['height'] = Inches(height)
    return slide.shapes.add_picture(path, **kwargs)


def title_slide(title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BLUE)
    add_textbox(slide, 0.5, 2.0, 12, 1.5, title, font_size=36, bold=True,
                color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, 0.5, 3.8, 12, 1.0, subtitle, font_size=18,
                    color=WHITE, alignment=PP_ALIGN.CENTER)
    return slide


def section_slide(title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MED_BLUE)
    add_textbox(slide, 0.5, 2.5, 12, 1.5, title, font_size=32, bold=True,
                color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, 0.5, 4.0, 12, 0.8, subtitle, font_size=16,
                    color=WHITE, alignment=PP_ALIGN.CENTER)
    return slide


def content_slide(title, description_lines, images=None, img_layout='full'):
    """
    images: list of (path, left, top, width, height) or (path, left, top, width)
    img_layout: 'full', 'left', 'right', 'two_col', 'stacked'
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title bar
    add_textbox(slide, 0.3, 0.1, 12.5, 0.6, title, font_size=22, bold=True,
                color=DARK_BLUE)

    # Description text
    if description_lines:
        if img_layout == 'full':
            add_multiline(slide, 0.3, 0.7, 12.5, 1.2, description_lines,
                          font_size=11, color=DARK_GRAY)
        elif img_layout in ('left', 'two_col', 'stacked'):
            add_multiline(slide, 0.3, 0.7, 12.5, 1.0, description_lines,
                          font_size=11, color=DARK_GRAY)

    # Images
    if images:
        for img_spec in images:
            path = img_spec[0]
            left = img_spec[1]
            top = img_spec[2]
            w = img_spec[3] if len(img_spec) > 3 else None
            h = img_spec[4] if len(img_spec) > 4 else None
            add_image(slide, path, left, top, w, h)

    return slide


def text_slide(title, lines, font_size=14):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.3, 0.1, 12.5, 0.6, title, font_size=24, bold=True,
                color=DARK_BLUE)
    add_multiline(slide, 0.5, 0.9, 12, 6.0, lines, font_size=font_size,
                  color=DARK_GRAY)
    return slide


# =============================================================================
# BUILD PRESENTATION
# =============================================================================
print("Building presentation...")

# ---- SLIDE 1: Title ----
title_slide(
    "Neural Signatures of Foraging Strategy",
    "LHA and RSP dynamics during pot visits, transitions, and reward\n"
    "Single-probe Neuropixels 2.0 | Mouse01 | Sessions 2, 4 (Fed), 6, 8 (Fasted)"
)

# ---- SLIDE 2: Task structure ----
content_slide(
    "Behavioral Overview: Four Foraging Sessions",
    [
        ("Fed mice (S2, S4) discover food after 767-937s; fasted mice (S6, S8) within 30s.", True, DARK_BLUE),
        "Arena has Pot-2 (food visible in exploration) and Pot-4 (food hidden in sand).",
        "Discovery = first feeding event. Pre-discovery: mouse explores; post-discovery: mouse exploits.",
        "Fed sessions have 32-39 pre-discovery visits; fasted sessions have only 2-3.",
    ],
    [(FIG / "foraging_cross_session_behavioral.png", 0.3, 2.2, 12.5, 5.0)],
)

# ---- SLIDE 3: Behavioral summary table ----
content_slide(
    "Behavioral and Neural Summary: All Sessions",
    [
        "Four-panel summary: (A) Behavioral, (B) Learning curves, (C) Pre vs post, (D) Key findings.",
        ("Fed sessions have extensive pre-discovery data; fasted sessions lack it entirely.", False, DARK_GRAY),
    ],
    [(FIG / "foraging_key_findings_table.png", 0.5, 1.8, 12, 5.5)],
)

# ---- SLIDE 4: Metrics primer ----
text_slide(
    "Neural Metrics: What We Measure",
    [
        ("Four metrics extracted at each pot visit from GRU-ODE latent dynamics:", True, DARK_BLUE),
        "",
        ("FR (Population Firing Rate)", True),
        "Mean firing rate across all good units in the region. Direct measure of population activity level.",
        "",
        ("PC1 (First Principal Component of Latent State)", True),
        "Dominant axis of variation in the 32-dimensional GRU-ODE hidden state. Captures the largest source of neural state change.",
        "",
        ("Flow Speed", True),
        "Magnitude of dh/dt from the ODE. How fast the latent state is evolving. High flow = rapid neural dynamics.",
        "",
        ("Gate Value", True),
        "Mean update gate from GRU-ODE. High gate = dynamics frozen (state preserved). Low gate = dynamics active (state changing).",
        "",
        ("Shuffle control: randomly reassign pot labels among all pre/post-discovery visits (10,000x), recompute correlation. Tests whether trends are pot-specific vs session-wide drift.", False, RED),
    ],
    font_size=13,
)

# ===== SECTION 1: POT-2 =====
section_slide("Section 1: Pot-2 (Obvious Food Pot)",
              "Pre- and post-discovery neural signatures")

# ---- SLIDE 6: Events timeline ----
content_slide(
    "Event Timelines: Fed Sessions (S2 and S4)",
    [
        "Top: Pot occupancy. Middle: Digging and feeding events. Bottom: Home vs arena.",
        ("S2 discovers food at 767s. S4 first digs at Pot-2 (783s), discovers food at Pot-4 at 937s.", True, DARK_BLUE),
    ],
    [
        (FIG / "foraging_events_s2.png", 0.2, 1.7, 12.8, 2.7),
        (FIG / "foraging_events_s4.png", 0.2, 4.5, 12.8, 2.7),
    ],
)

# ---- SLIDE 7: Pre-disc learning curves fed ----
content_slide(
    "Pre-Discovery Learning Curves: Fed Sessions",
    [
        "Each dot = one pot visit. Correlation of metric vs visit number (Pearson r).",
        ("S2 LHA PC1 at Pot-2: r=0.80, p=0.056. S4 LHA PC1 at Pot-2: r=0.76, p=0.0004.", True, DARK_BLUE),
        ("But do these survive shuffle control? See next slide.", False, RED),
    ],
    [
        (FIG / "foraging_neural_across_visits_s2.png", 0.2, 2.0, 6.3, 5.0),
        (FIG / "foraging_neural_across_visits_s4.png", 6.7, 2.0, 6.3, 5.0),
    ],
)

# ---- SLIDE 8: Pre-disc fasted ----
content_slide(
    "Pre-Discovery: Fasted Sessions Skip Learning Entirely",
    [
        ("S6 and S8 have ZERO pre-discovery Pot-2 visits. Only 1-2 Pot-4 visits before finding food.", True, RED),
        "Fasted mice discover food within 30-34 seconds. The 700-900s exploration period in fed mice does not exist.",
        "No learning curve analysis possible for fasted sessions.",
    ],
    [
        (FIG / "foraging_neural_across_visits_s6.png", 0.2, 2.0, 6.3, 5.0),
        (FIG / "foraging_neural_across_visits_s8.png", 6.7, 2.0, 6.3, 5.0),
    ],
)

# ---- SLIDE 9: Shuffle control PC1 ----
content_slide(
    "Shuffle Control: LHA PC1 Pot-2 Ramp Does NOT Survive",
    [
        ("S4: p_shuffle=0.115 (ns). Strong session-wide drift (all-visits r=0.72, p<0.0001) explains the ramp.", True, RED),
        ("S2: p_shuffle=0.055 (borderline). No drift, but underpowered (n=6).", False, DARK_GRAY),
        ("Right panels: ALL pot visits trend upward in S4 — drift, not pot-specific learning.", False, DARK_GRAY),
    ],
    [(FIG / "foraging_shuffle_control_pc1.png", 0.2, 2.0, 12.8, 5.0)],
)

# ---- SLIDE 10: What survives pre-disc ----
content_slide(
    "Pre-Discovery Shuffle: What Survives",
    [
        ("Only S2 LHA Pot-4 signals survive: Flow (p=0.020*) and Gate (p=0.022*).", True, GREEN),
        "No Pot-2 metrics survive shuffle in any session. The pre-discovery story is about Pot-4, not Pot-2.",
        "Dots outside gray bars = observed r exceeds 95% CI of shuffle distribution.",
    ],
    [(FIG / "foraging_shuffle_control_all_metrics.png", 0.2, 2.0, 12.8, 5.0)],
)

# ---- SLIDE 11: Pre vs post stats ----
content_slide(
    "Pre vs Post-Discovery: LHA Firing Rate Drops at Pot-2",
    [
        ("S4: LHA FR at Pot-2 significantly lower post-discovery (1.80 → 1.31, p=0.030*).", True, GREEN),
        ("S2: Same trend but not significant (1.88 → 1.63, p=0.066).", False, DARK_GRAY),
        "This is a between-condition comparison (pre vs post), not a temporal trend — drift does not apply.",
    ],
    [
        (FIG / "foraging_pre_vs_post_stats_s4.png", 0.2, 2.2, 6.3, 5.0),
        (FIG / "foraging_pre_vs_post_stats_s2.png", 6.7, 2.2, 6.3, 5.0),
    ],
)

# ---- SLIDE 12: Post-disc shuffle Pot-2 ----
content_slide(
    "Post-Discovery Pot-2: S6 Shows LHA-RSP Dissociation",
    [
        ("S6 (Fasted): LHA Pot-2 FR ↓ (p=0.050*), PC1 ↓ (p=0.023*) — pot-specific, no drift.", True, GREEN),
        ("S6 RSP Pot-2: PC1 ↑ (p=0.031*), Flow ↑ (p=0.014*) — opposite direction from LHA.", True, MED_BLUE),
        "Coordinated but opposing LHA-RSP modulation at Pot-2 post-discovery. Cleanest signal in dataset.",
    ],
    [(FIG / "foraging_shuffle_post_disc_all_metrics.png", 0.2, 2.2, 12.8, 5.0)],
)

# ===== SECTION 2: POT-4 =====
section_slide("Section 2: Pot-4 (Hidden Food Pot)",
              "Pre- and post-discovery neural signatures")

# ---- SLIDE 14: Pre-disc Pot-4 ----
content_slide(
    "Pre-Discovery Pot-4: The Only Genuine Learning Signal",
    [
        ("S2 LHA Pot-4: Flow decreases (r=-0.52, p_shuffle=0.020*) and Gate increases (r=+0.64, p_shuffle=0.022*).", True, GREEN),
        "These survive label shuffle — genuinely pot-specific, not drift.",
        "Interpretation: LHA dynamics slow down and freeze at the hidden food pot before discovery.",
        ("S4 has only 6 pre-disc Pot-4 visits — nothing survives.", False, DARK_GRAY),
    ],
    [(FIG / "foraging_shuffle_control_all_metrics.png", 0.2, 2.5, 12.8, 4.5)],
)

# ---- SLIDE 15: Post-disc Pot-4 gate ----
content_slide(
    "Post-Discovery Pot-4: LHA Gate Survives But Flips Direction",
    [
        ("S2: LHA Pot-4 Gate decreasing (r=-0.44, p_shuffle=0.027*). S4: increasing (r=+0.60, p_shuffle=0.020*).", True, GOLD),
        "Both survive shuffle but show opposite trends. S4 first dug at Pot-2 — may explain the sign flip.",
        "S6/S8 fasted: no Pot-4 signals survive shuffle.",
    ],
    [(FIG / "foraging_shuffle_post_disc_time.png", 0.2, 2.0, 12.8, 5.2)],
)

# ---- SLIDE 16: Fasted pre/post ----
content_slide(
    "Fasted Sessions: No Pre vs Post Changes at Pot-4",
    [
        "S6: 2 pre vs 23 post Pot-4 visits — Mann-Whitney all p≥0.53. Too few pre-disc visits.",
        "S8: 1 pre vs 56 post — cannot run statistics. RSP drift dominates (r=0.40, p=0.0002).",
        ("Fasted mice show no detectable pre/post neural shift at the food pot.", True, RED),
    ],
    [
        (FIG / "foraging_pre_vs_post_stats_s6.png", 0.2, 2.3, 6.3, 4.8),
        (FIG / "foraging_pre_vs_post_stats_s8.png", 6.7, 2.3, 6.3, 4.8),
    ],
)

# ---- SLIDE 17: S8 drift example ----
content_slide(
    "S8 RSP: Session-Wide Drift Dominates",
    [
        ("S8 RSP shows strong drift: FR r=+0.40 (p=0.0002), PC1 r=+0.40 (p=0.0002) over all post-disc visits.", True, RED),
        "This is the clearest example of drift overwhelming pot-specific signals.",
        "No post-discovery metrics survive shuffle in S8 for either region.",
    ],
    [(FIG / "foraging_continuous_s8_rsp.png", 0.5, 2.0, 12, 5.0)],
)

# ===== SECTION 3: TRANSITIONS =====
section_slide("Section 3: P2 ↔ P4 Transitions",
              "Within-excursion neural state changes")

# ---- SLIDE 19: Transition paired ----
content_slide(
    "Within-Excursion: P2 vs P4 Neural State (16 Transitions)",
    [
        "Lines connect P2 and P4 values within the same excursion. Black = pooled mean.",
        ("LHA Gate: P4 > P2, Wilcoxon p=0.044* — 12 of 16 transitions show higher gate at Pot-4.", True, GREEN),
        ("RSP PC1: P4 < P2, trending (p=0.074) — 11 of 16 transitions.", False, DARK_GRAY),
    ],
    [(FIG / "foraging_transition_paired.png", 0.2, 2.0, 12.8, 5.2)],
)

# ---- SLIDE 20: Transition deltas ----
content_slide(
    "Transition Deltas by Session (P4 minus P2)",
    [
        "Each dot = one transition. Horizontal bars = session mean. Blue = fed, red = fasted.",
        ("LHA Gate is the only metric with consistent P4 > P2 difference across sessions.", True, GREEN),
        "Fed vs fasted: no significant differences in any metric (all MWU p>0.16).",
    ],
    [(FIG / "foraging_transition_deltas.png", 0.2, 2.0, 12.8, 5.2)],
)

# ---- SLIDE 21: Heatmap ----
content_slide(
    "Transition Heatmap: Per-Session and Pooled",
    [
        "Green = P4 higher, Red = P2 higher. Bold = session-level Wilcoxon p<0.05.",
        ("S4 drives the signal: LHA Gate 6/6 positive (p=0.031*), RSP FR 6/6 positive (p=0.031*).", True, GREEN),
        "S6 RSP FR flips direction (P4 < P2) vs fed sessions.",
    ],
    [(FIG / "foraging_transition_heatmap.png", 0.5, 2.0, 12, 5.0)],
)

# ---- SLIDE 22: S4 transition dynamics ----
content_slide(
    "S4 Transition Dynamics: Strongest Session",
    [
        "S4 has 6 post-discovery P2↔P4 transitions. Rows: FR, PC1, Flow, Gate, Divergence.",
        ("All 6 transitions show LHA Gate higher during Pot-4 occupancy (green) vs Pot-2 (pink).", True, GREEN),
    ],
    [
        (FIG / "foraging_transition_dynamics_s4_lha.png", 0.2, 1.8, 12.8, 5.5),
    ],
)

# ---- SLIDE 23: Cross-session transitions ----
content_slide(
    "Cross-Session Transition Summary",
    [
        "Left: transition counts by phase. Middle: P2→P4 vs P4→P2 direction. Right: gap time between pots.",
        "Fasted mice have shorter gap times (~3s vs ~18s in fed).",
        "S4 has the most transitions (12 total, 7 post-discovery).",
    ],
    [(FIG / "foraging_cross_session_transitions.png", 0.3, 1.8, 12.5, 5.2)],
)

# ===== SECTION 4: REWARD ======
section_slide("Reward Onset Responses",
              "Neural dynamics at the moment of food discovery")

# ---- SLIDE 25: Reward fed ----
content_slide(
    "Reward Onset: Fed Sessions (S2, S4)",
    [
        "Peri-reward neural response at first feeding event. Time 0 = feed onset.",
        "LHA FR rises post-reward in both sessions. LHA PC1 shows sharp transients.",
    ],
    [
        (FIG / "foraging_reward_s2.png", 0.2, 1.8, 6.3, 5.2),
        (FIG / "foraging_reward_s4.png", 6.7, 1.8, 6.3, 5.2),
    ],
)

# ---- SLIDE 26: Reward fasted ----
content_slide(
    "Reward Onset: Fasted Sessions (S6, S8)",
    [
        "S6 RSP shows dramatic FR drop at reward onset. Less consistent than fed sessions.",
        "LHA PC1 transients at feed onset are present in all sessions.",
    ],
    [
        (FIG / "foraging_reward_s6.png", 0.2, 1.8, 6.3, 5.2),
        (FIG / "foraging_reward_s8.png", 6.7, 1.8, 6.3, 5.2),
    ],
)

# ---- SLIDE 27: Cross-session reward ----
content_slide(
    "Cross-Session Reward Comparison",
    [
        "Fed (blue) vs fasted (red) reward onset responses overlaid.",
        "Fed sessions show more consistent LHA FR increase post-reward.",
    ],
    [(FIG / "foraging_cross_session_reward.png", 0.5, 1.8, 12, 5.2)],
)

# ===== SECTION 5: SYNTHESIS =====
section_slide("Synthesis",
              "What survived, what didn't, and what it means")

# ---- SLIDE 29: Summary table ----
text_slide(
    "What Survived Shuffle Controls",
    [
        ("All results that survive 10,000-permutation label shuffle (p < 0.05):", True, DARK_BLUE),
        "",
        ("PRE-DISCOVERY (Fed only — fasted have insufficient data):", True, MED_BLUE),
        "  S2 LHA Pot-4 Flow: r=-0.52, p_shuffle=0.020* — dynamics slow at hidden food",
        "  S2 LHA Pot-4 Gate: r=+0.64, p_shuffle=0.022* — activation diversifies at hidden food",
        ("  LHA PC1 Pot-2 ramp does NOT survive (S4 driven by session drift)", False, RED),
        "",
        ("POST-DISCOVERY:", True, MED_BLUE),
        "  S6 LHA Pot-2 FR: r=-0.86, p_shuffle=0.050* — firing drops at empty pot",
        "  S6 LHA Pot-2 PC1: r=-0.90, p_shuffle=0.023* — latent state shifts away",
        "  S6 RSP Pot-2 PC1: r=+0.89, p_shuffle=0.031* — RSP shifts opposite to LHA",
        "  S6 RSP Pot-2 Flow: r=+0.92, p_shuffle=0.014* — RSP dynamics accelerate",
        "  S2 LHA Pot-4 Gate: r=-0.44, p_shuffle=0.027*",
        "  S4 LHA Pot-4 Gate: r=+0.60, p_shuffle=0.020* (opposite direction to S2)",
        "",
        ("WITHIN-EXCURSION TRANSITIONS (pooled, n=16):", True, MED_BLUE),
        "  LHA Gate: P4 > P2, Wilcoxon p=0.044* (12/16 transitions consistent)",
        "  S4 alone: LHA Gate p=0.031* (6/6), RSP FR p=0.031* (6/6)",
        "",
        ("PRE vs POST COMPARISON:", True, MED_BLUE),
        "  S4 LHA FR at Pot-2: drops post-discovery (p=0.030*) — value update",
    ],
    font_size=12,
)

# ---- SLIDE 30: The LHA Gate story ----
content_slide(
    "The Consistent Signal: LHA Gate Differentiates Pots",
    [
        ("LHA Gate is the most robust neural signal across analyses:", True, DARK_BLUE),
        ("1. Pre-discovery: increases at Pot-4 over visits (S2, p_shuffle=0.022*)", False, GREEN),
        ("2. Within-excursion: higher at Pot-4 than Pot-2 (pooled p=0.044*, 12/16 transitions)", False, GREEN),
        ("3. Post-discovery: changes over time at Pot-4 in both fed sessions (both p<0.03*)", False, GREEN),
        ("Higher gate = more frozen dynamics = state preservation at the food pot.", True, DARK_BLUE),
    ],
    [
        (FIG / "foraging_transition_paired.png", 0.2, 3.0, 6.3, 4.2),
        (FIG / "foraging_shuffle_control_all_metrics.png", 6.7, 3.0, 6.3, 4.2),
    ],
)

# ---- SLIDE 31: Fed vs fasted ----
content_slide(
    "Fed vs Fasted: Different Strategies, Shared Pot Discrimination",
    [
        "Fed mice explore for 700-900s before discovery; fasted mice find food within 30s.",
        ("Both states show LHA Gate differentiation at Pot-4 in transitions (no fed vs fasted difference, p>0.16).", True, DARK_BLUE),
        "S6 fasted shows the cleanest post-discovery signal: LHA-RSP dissociation at Pot-2.",
    ],
    [(FIG / "foraging_cross_session_fed_vs_fasted.png", 0.5, 2.0, 12, 5.0)],
)

# ---- SLIDE 32: Take-home ----
text_slide(
    "Take-Home Messages",
    [
        ("1. LHA Gate is the most consistent neural signal", True, DARK_BLUE),
        "   It differentiates Pot-4 (food) from Pot-2 (empty) within excursions (p=0.044),",
        "   shows pot-specific pre-discovery trends (p=0.022), and is robust to shuffle control.",
        "   Higher gate = frozen dynamics = LHA preserves its neural state at the food pot.",
        "",
        ("2. Pre-discovery 'learning curves' at Pot-2 are largely drift", True, DARK_BLUE),
        "   The LHA PC1 ramp (S4 p=0.0004 parametric) does NOT survive shuffle (p=0.115).",
        "   Session-wide temporal drift (r=0.72) explains the apparent learning signal.",
        "   Always run shuffle controls before interpreting temporal neural trends.",
        "",
        ("3. S6 fasted reveals genuine LHA-RSP dissociation post-discovery", True, DARK_BLUE),
        "   LHA metrics decrease and RSP metrics increase at Pot-2 — opposing modulation.",
        "   All four signals survive shuffle with no session-wide drift. Pot-specific.",
        "",
        ("4. Fasted mice skip learning entirely", True, DARK_BLUE),
        "   Zero pre-discovery Pot-2 visits, food found within 30s.",
        "   The 700-900s exploration in fed mice simply does not exist in fasted state.",
    ],
    font_size=13,
)

# ---- SLIDE 33: Caveats ----
text_slide(
    "Caveats and Next Steps",
    [
        ("Limitations:", True, RED),
        "  • N=1 mouse, N=2 sessions per metabolic state — need replication",
        "  • S4 LHA Pot-4 Gate flips direction post-discovery (unexplained)",
        "  • Session drift is pervasive — some sessions (S4 pre, S8 post) show strong global trends",
        "  • Small n per session limits power for within-session shuffle tests",
        "  • Fasted sessions have too few pre-discovery visits for any statistical analysis",
        "",
        ("Next steps:", True, GREEN),
        "  • More sessions / mice to replicate LHA Gate finding",
        "  • Dual-probe recordings (ACA + LHA) — does ACA contribute to the Gate signal?",
        "  • Neuropixels-Opto causal perturbation — is LHA Gate necessary for pot discrimination?",
        "  • Cross-session unit tracking (UnitMatch) — do the same LHA neurons drive the Gate effect?",
        "  • Control for movement speed / acceleration at pot arrival",
        "  • Extend to exploration sessions to compare pot representations with/without food",
    ],
    font_size=13,
)

# ---- SLIDE 34: End ----
title_slide(
    "Thank You",
    "Data, scripts, and figures available in the NPX Analysis Repo"
)

# =============================================================================
# SAVE
# =============================================================================
out_path = "data/foraging_neural_signatures_presentation.pptx"
prs.save(out_path)
print(f"\nSaved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
