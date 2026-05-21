"""
Summary PowerPoint: Foraging Neural Signatures (concise version).

~12 slides focused on key findings:
1. Pot-2 (pre/post discovery)
2. Pot-4 (pre/post discovery)
3. P2<->P4 transitions
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FIG = Path("figures")
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x2E, 0x75, 0xB6)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
GOLD = RGBColor(0xF3, 0x9C, 0x12)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=14,
             bold=False, color=BLACK, align=PP_ALIGN.LEFT, name='Calibri'):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return tb


def add_lines(slide, left, top, width, height, lines, size=12,
              color=BLACK, name='Calibri', spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, bold, col = line, False, color
        else:
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            col = line[2] if len(line) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = col
        p.font.name = name
        p.space_after = Pt(size * spacing * 0.5)
    return tb


def add_img(slide, path, left, top, width=None, height=None):
    path = str(path)
    if not os.path.exists(path):
        print(f"  WARNING: {path} not found, skipping")
        return None
    kw = {'left': Inches(left), 'top': Inches(top)}
    if width:
        kw['width'] = Inches(width)
    if height:
        kw['height'] = Inches(height)
    return slide.shapes.add_picture(path, **kw)


# ─────────────────────────────────────────────────────────
# SLIDE 1  — Title
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BLUE)
add_text(s, 0.5, 1.8, 12.3, 1.5,
         "Neural Signatures of Foraging Strategy",
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 3.5, 12.3, 0.6,
         "LHA & RSP latent dynamics during pot visits and transitions",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.3, 12.3, 0.6,
         "Single-probe Neuropixels 2.0  |  Mouse01  |  S2, S4 (Fed)  S6, S8 (Fasted)",
         size=14, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
add_lines(s, 0.5, 5.5, 12.3, 1.0, [
    ("Metrics: Population FR, Latent PC1, Trajectory Flow Speed, GRU-ODE Gate", False, RGBColor(0x99, 0xAA, 0xCC)),
    ("All temporal trends validated with 10,000-permutation shuffle controls", False, RGBColor(0x99, 0xAA, 0xCC)),
], size=12, color=RGBColor(0x99, 0xAA, 0xCC))

# ─────────────────────────────────────────────────────────
# SLIDE 2  — Behavioral overview (brief)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.15, 12.5, 0.5,
         "Behavioral Overview", size=24, bold=True, color=DARK_BLUE)
add_lines(s, 0.4, 0.75, 5.5, 5.5, [
    ("Paradigm", True, DARK_BLUE),
    "Exploration sessions: visible food at Pot-2",
    "Foraging sessions: food hidden in sand (Pot-3 or Pot-4)",
    "Mouse must learn Pot-2 is now empty and find new food location",
    "",
    ("Fed mice (S2, S4)", True, MED_BLUE),
    "  Discovery at 767s (S2) and 937s (S4)",
    "  S2: 6 pre-disc Pot-2 visits, 16 Pot-4 visits",
    "  S4: 17 pre-disc Pot-2 visits, 6 Pot-4 visits (digs at P2 first)",
    "",
    ("Fasted mice (S6, S8)", True, RED),
    "  Discovery at 34s (S6) and 30s (S8)",
    "  0 pre-disc Pot-2 visits — skip directly to food pot",
    "  All learning signals are post-discovery only",
], size=12, color=DARK_GRAY)
add_img(s, FIG / "foraging_cross_session_behavioral.png", 6.2, 0.8, width=6.8)


# ─────────────────────────────────────────────────────────
# SLIDE 3  — Section: Pot-2
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, MED_BLUE)
add_text(s, 0.5, 2.5, 12.3, 1.2,
         "Section 1: Pot-2 — The Previously-Rewarded Pot",
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.0, 12.3, 0.8,
         "How does the brain update value of a pot that no longer has food?",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────
# SLIDES 4-7  — Pot-2: continuous LHA traces, one per session
# ─────────────────────────────────────────────────────────
continuous_pot2 = [
    ("S2", "Fed", "foraging_continuous_s2_lha.png",
     ["6 pre-discovery Pot-2 visits (red shading), 16 Pot-4 visits (green).",
      "Discovery at 767s. LHA FR trending lower at Pot-2 post-discovery (p=0.066).",
      ("Flow, Gate, PC1, FR tracked continuously across session.", True, DARK_GRAY)]),
    ("S4", "Fed", "foraging_continuous_s4_lha.png",
     ["17 pre-discovery Pot-2 visits (red shading) — most in any session.",
      "Mouse digs at Pot-2 first (783s) before finding food at Pot-4 (937s).",
      ("LHA FR at Pot-2 drops significantly post-discovery (p=0.030*).", True, DARK_GRAY)]),
    ("S6", "Fasted", "foraging_continuous_s6_lha.png",
     ["0 pre-discovery Pot-2 visits. Food found at Pot-4 in just 34s.",
      "Post-discovery: 6 Pot-2 visits show coordinated LHA-RSP updating.",
      ("Strongest pot-specific signals in the dataset (all survive shuffle).", True, DARK_GRAY)]),
    ("S8", "Fasted", "foraging_continuous_s8_lha.png",
     ["0 pre-discovery Pot-2 visits. Food found at Pot-4 in just 30s.",
      "Post-discovery: 8 Pot-2 visits, 56 Pot-4 visits (highest revisit count).",
      ("RSP drift dominates — no pot-specific signals survive shuffle.", True, DARK_GRAY)]),
]
for sess, state, fig, desc in continuous_pot2:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_text(s, 0.3, 0.1, 12.5, 0.5,
             f"Pot-2: LHA Dynamics ({sess}, {state})",
             size=22, bold=True, color=DARK_BLUE)
    add_lines(s, 0.3, 0.6, 5.0, 1.5, desc, size=11, color=DARK_GRAY)
    add_img(s, FIG / fig, 0.2, 1.5, width=12.8, height=5.8)


# ─────────────────────────────────────────────────────────
# SLIDE 8  — Pot-2: pre vs post feed stats (S4 significant)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "Pot-2: Pre vs Post-Discovery Neural State",
         size=22, bold=True, color=DARK_BLUE)
# S4 on left, S2 on right
add_img(s, FIG / "foraging_pre_vs_post_stats_s4.png", 0.1, 1.8, width=6.5)
add_img(s, FIG / "foraging_pre_vs_post_stats_s2.png", 6.6, 1.8, width=6.5)
add_lines(s, 0.3, 0.65, 12.5, 1.2, [
    ("S4 (left): LHA FR at Pot-2 drops significantly post-discovery (pre=1.80, post=1.31, p=0.030*)", True, DARK_GRAY),
    "S2 (right): Same trend, borderline (pre=1.88, post=1.63, p=0.066)",
    "This is a direct pre/post comparison — immune to session-wide drift concerns.",
], size=11, color=DARK_GRAY)
# Labels
add_text(s, 2.0, 1.35, 3.0, 0.4, "S4 (Fed)", size=14, bold=True, color=MED_BLUE,
         align=PP_ALIGN.CENTER)
add_text(s, 8.5, 1.35, 3.0, 0.4, "S2 (Fed)", size=14, bold=True, color=MED_BLUE,
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────
# SLIDE 6  — Pot-2: S6 post-disc shuffle (strongest signal)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "Pot-2: Post-Discovery Learning — Shuffle-Validated (S6, Fasted)",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.6, 5.5, 2.0, [
    ("S6 shows the strongest pot-specific signals in the dataset:", True, DARK_GRAY),
    "",
    ("LHA Pot-2 FR:  r=-0.86, p_shuffle=0.050*  (decreasing)", False, MED_BLUE),
    ("LHA Pot-2 PC1: r=-0.90, p_shuffle=0.023*  (decreasing)", False, MED_BLUE),
    ("RSP Pot-2 PC1: r=+0.89, p_shuffle=0.031*  (increasing)", False, RED),
    ("RSP Pot-2 Flow: r=+0.92, p_shuffle=0.014* (increasing)", False, RED),
    "",
    "LHA and RSP move in opposite directions at the empty pot.",
    "No session-wide drift (all-visits r = -0.02). Genuinely pot-specific.",
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_shuffle_post_disc_all_metrics.png", 5.8, 0.6, width=7.3, height=6.5)


# ─────────────────────────────────────────────────────────
# SLIDE 7  — Section: Pot-4
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, MED_BLUE)
add_text(s, 0.5, 2.5, 12.3, 1.2,
         "Section 2: Pot-4 — The Hidden Food Pot",
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.0, 12.3, 0.8,
         "How does the brain respond to reward discovery and subsequent revisits?",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────
# SLIDES — Pot-4: peri-reward per session + cross-session overlay
# ─────────────────────────────────────────────────────────
peri_reward = [
    ("S2", "Fed", "foraging_reward_s2.png",
     ["Discovery at 767s. LHA FR rises post-reward; LHA PC1 shows sharp transient.",
      ("6 panels: FR, PC1-3, Flow & Gate for LHA (top) and RSP (bottom).", True, DARK_GRAY)]),
    ("S4", "Fed", "foraging_reward_s4.png",
     ["Discovery at 937s. LHA FR ramps up after feeding onset.",
      ("Similar pattern to S2 — consistent fed-state reward response.", True, DARK_GRAY)]),
    ("S6", "Fasted", "foraging_reward_s6.png",
     ["Discovery at 34s. RSP shows dramatic FR drop at reward onset.",
      ("LHA PC1 transient present. Distinct from fed-session pattern.", True, DARK_GRAY)]),
    ("S8", "Fasted", "foraging_reward_s8.png",
     ["Discovery at 30s. Fastest discovery across all sessions.",
      ("Less consistent pattern — may reflect very rapid state transition.", True, DARK_GRAY)]),
]
for sess, state, fig, desc in peri_reward:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_text(s, 0.3, 0.1, 12.5, 0.5,
             f"Pot-4: Peri-Reward Response ({sess}, {state})",
             size=22, bold=True, color=DARK_BLUE)
    add_lines(s, 0.3, 0.55, 12.5, 0.8, desc, size=11, color=DARK_GRAY)
    add_img(s, FIG / fig, 0.2, 1.3, width=12.8, height=5.9)

# Cross-session overlay
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "Pot-4: Peri-Reward Overlay — Fed vs Fasted",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.65, 12.5, 0.9, [
    "All 4 sessions aligned to first feeding onset (t=0). LHA FR and PC1 (top), RSP FR and PC1 (bottom).",
    ("Fed (blue): LHA FR ramps up post-reward. Fasted (red): S6 RSP shows dramatic FR drop.", True, DARK_GRAY),
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_cross_session_reward.png", 0.5, 1.6, width=12.3, height=5.6)


# ─────────────────────────────────────────────────────────
# SLIDE 10  — Section: Transitions
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, MED_BLUE)
add_text(s, 0.5, 2.5, 12.3, 1.2,
         "Section 3: Pot-2 to Pot-4 Transitions",
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.0, 12.3, 0.8,
         "Within-excursion comparison: how does the neural state change\n"
         "as the mouse moves from the empty pot to the food pot?",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────
# SLIDES — Transitions: traces per session (LHA)
# ─────────────────────────────────────────────────────────
transition_sessions = [
    ("S2", "Fed", "foraging_transition_dynamics_s2_lha.png",
     "5 transitions (3 pre, 2 post-discovery). Each column = one excursion with both P2 and P4."),
    ("S4", "Fed", "foraging_transition_dynamics_s4_lha.png",
     "12 transitions (5 pre, 7 post). LHA Gate consistently higher at P4 — 6/6 post-disc (p=0.031*)."),
    ("S6", "Fasted", "foraging_transition_dynamics_s6_lha.png",
     "5 transitions (0 pre, 5 post). RSP FR flips direction (P4 < P2) vs fed sessions."),
    ("S8", "Fasted", "foraging_transition_dynamics_s8_lha.png",
     "4 transitions (0 pre, 4 post). All P2->P4 direction: consistent LHA FR drop at P4."),
]
for sess, state, fig, desc in transition_sessions:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_text(s, 0.3, 0.1, 12.5, 0.5,
             f"P2-P4 Transition Traces ({sess}, {state} — LHA)",
             size=22, bold=True, color=DARK_BLUE)
    add_lines(s, 0.3, 0.55, 12.5, 0.8, [
        desc,
        ("Red shading = Pot-2, green = Pot-4. Rows: FR, PC1, Flow, Gate, Divergence.", True, DARK_GRAY),
    ], size=11, color=DARK_GRAY)
    add_img(s, FIG / fig, 0.2, 1.35, width=12.8, height=6.0)


# ─────────────────────────────────────────────────────────
# SLIDE — Transitions: paired plot
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "P2-P4 Transitions: Paired Comparison Across All Sessions",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.6, 12.5, 0.8, [
    "16 post-discovery transitions across S2/S4/S6/S8. Lines connect P2 and P4 values within same excursion.",
    ("LHA Gate: P4 > P2 in 12/16 transitions (Wilcoxon p=0.044*). S4 drives: 6/6 positive (p=0.031*).", True, DARK_GRAY),
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_transition_paired.png", 0.1, 1.5, width=12.8, height=5.7)


# ─────────────────────────────────────────────────────────
# SLIDE 13  — Transitions: heatmap summary
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "P2-P4 Transition Deltas by Session (Heatmap)",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.6, 12.5, 0.8, [
    "Mean delta (Pot-4 minus Pot-2) per session. Green = P4 higher, Red = P2 higher. Bold = session Wilcoxon p<0.05.",
    ("LHA Gate consistently positive (P4 > P2). RSP PC1 trending negative (P4 < P2). RSP FR flips in fasted S6.", True, DARK_GRAY),
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_transition_heatmap.png", 0.3, 1.5, width=12.5, height=5.5)


# ─────────────────────────────────────────────────────────
# SLIDE — Shuffle-validated signals for Pot-2 and Pot-4
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "What Survived Shuffle Control? (Pot-2 & Pot-4, Post-Discovery)",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.7, 12.5, 0.8, [
    "10,000-permutation label shuffle: randomly reassign pot labels among all post-disc visits, recompute Pearson r.",
    ("Signals below survived — observed r falls outside the 95% CI of the shuffle distribution.", True, DARK_GRAY),
], size=11, color=DARK_GRAY)

# Pot-2 table
add_text(s, 0.4, 1.6, 5.0, 0.4, "Pot-2 (empty pot) — post-discovery trends",
         size=16, bold=True, color=MED_BLUE)
add_lines(s, 0.5, 2.1, 6.0, 3.5, [
    ("S6 (Fasted) — coordinated LHA-RSP dissociation:", True, DARK_GRAY),
    "",
    ("  LHA FR      r = -0.86   p_shuffle = 0.050*   decreasing", False, MED_BLUE),
    ("  LHA PC1     r = -0.90   p_shuffle = 0.023*   decreasing", False, MED_BLUE),
    ("  RSP PC1     r = +0.89   p_shuffle = 0.031*   increasing", False, RED),
    ("  RSP Flow    r = +0.92   p_shuffle = 0.014*   increasing", False, RED),
    "",
    "No session-wide drift (all-visits r = -0.02).",
    "LHA and RSP move in opposite directions at the empty pot.",
    "",
    ("S2, S4, S8: nothing at Pot-2 survives shuffle.", False, DARK_GRAY),
], size=12, color=DARK_GRAY)

# Pot-4 table
add_text(s, 6.8, 1.6, 5.5, 0.4, "Pot-4 (food pot) — post-discovery trends",
         size=16, bold=True, color=GREEN)
add_lines(s, 6.9, 2.1, 6.0, 3.5, [
    ("LHA Gate is the only metric that survives:", True, DARK_GRAY),
    "",
    ("  S2 LHA Gate  r = -0.44   p_shuffle = 0.027*   decreasing", False, MED_BLUE),
    ("  S4 LHA Gate  r = +0.60   p_shuffle = 0.020*   increasing", False, MED_BLUE),
    "",
    "Opposite directions between S2 and S4.",
    "May reflect different post-discovery revisit strategies.",
    "",
    ("S6, S8: nothing at Pot-4 survives shuffle.", False, DARK_GRAY),
    "",
    ("Within-excursion test (not shuffle):", True, DARK_GRAY),
    ("  LHA Gate: P4 > P2 in 12/16 transitions", False, DARK_GRAY),
    ("  Wilcoxon p = 0.044*", False, DARK_GRAY),
], size=12, color=DARK_GRAY)

# Bottom summary
add_lines(s, 0.4, 5.8, 12.5, 1.2, [
    ("Bottom line: LHA Gate is the most consistent pot-discriminating signal — survives shuffle at Pot-4,", True, DARK_BLUE),
    ("and distinguishes P4 from P2 within excursions. S6 Pot-2 LHA-RSP dissociation is the strongest single-session signal.", True, DARK_BLUE),
], size=13, color=DARK_BLUE)


# ─────────────────────────────────────────────────────────
# SLIDE — Transition shuffle control
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "P2-P4 Transition Shuffle Control (10,000 permutations)",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.6, 6.0, 1.0, [
    "For each permutation: randomly flip P2/P4 labels within each transition,",
    "recompute Wilcoxon statistic. Red line = observed statistic.",
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_transition_shuffle.png", 0.3, 1.5, width=7.0, height=3.0)

add_text(s, 7.5, 1.5, 5.5, 0.4, "Pooled (n=16 transitions)",
         size=16, bold=True, color=DARK_BLUE)
add_lines(s, 7.6, 2.0, 5.5, 2.5, [
    ("LHA Gate:  p_shuffle = 0.048*  SURVIVES", True, GREEN),
    ("RSP PC1:   p_shuffle = 0.075   trending", False, GOLD),
    "LHA FR:    p_shuffle = 0.98    ns",
    "LHA PC1:   p_shuffle = 0.82    ns",
    "LHA Flow:  p_shuffle = 0.48    ns",
    "RSP FR:    p_shuffle = 0.44    ns",
    "RSP Flow:  p_shuffle = 0.35    ns",
    "RSP Gate:  p_shuffle = 0.27    ns",
], size=11, color=DARK_GRAY)

add_text(s, 7.5, 4.5, 5.5, 0.4, "Per-session (shuffle-validated)",
         size=16, bold=True, color=DARK_BLUE)
add_lines(s, 7.6, 4.95, 5.5, 1.5, [
    ("S4 LHA Gate:  p_shuffle = 0.031*  (6/6 positive)", True, GREEN),
    ("S4 RSP FR:    p_shuffle = 0.034*  (6/6 positive)", True, GREEN),
    "S6, S8: nothing survives (n=4, insufficient power)",
    "S2: n=2, cannot test",
], size=11, color=DARK_GRAY)

add_lines(s, 0.3, 6.2, 12.5, 0.8, [
    ("LHA Gate is the only metric that survives both within-excursion Wilcoxon AND permutation shuffle.", True, DARK_BLUE),
    ("S4 drives the effect: 6/6 transitions show higher Gate at Pot-4 than Pot-2.", True, DARK_BLUE),
], size=13, color=DARK_BLUE)


# ─────────────────────────────────────────────────────────
# SLIDE — Section: Transit Dynamics
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, MED_BLUE)
add_text(s, 0.5, 2.5, 12.3, 1.2,
         "Section 4: Transit Dynamics — Decision Points",
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.0, 12.3, 0.8,
         "Full pot-to-pot neural trajectories: departure, transit, arrival.\n"
         "When does the brain 'decide' to leave?",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────
# SLIDE — Transit trajectory example (S4 LHA, richest session)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "Transit Trajectories: S4 (Fed) — LHA",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.55, 12.5, 0.8, [
    "8 post-discovery P2<->P4 transits. Each column = one transit. Rows: FR, PC1, Flow, Gate.",
    ("Dashed = departure, dotted = arrival. Red shading = Pot-2, green = Pot-4.", True, DARK_GRAY),
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_transit_traj_s4_lha.png", 0.2, 1.35, width=12.8, height=6.0)


# ─────────────────────────────────────────────────────────
# SLIDE — Transit trajectory example (S4 RSP)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "Transit Trajectories: S4 (Fed) — RSP",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.55, 12.5, 0.8, [
    "Same 8 transits as previous slide, now showing RSP dynamics.",
    ("RSP FR tends higher during P2->P4 transitions. Gate shows varied patterns.", True, DARK_GRAY),
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_transit_traj_s4_rsp.png", 0.2, 1.35, width=12.8, height=6.0)


# ─────────────────────────────────────────────────────────
# SLIDE — Change-point detection summary
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "Change-Point Detection: When Does the Neural State Shift?",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.6, 5.5, 1.0, [
    "CUSUM change-point detection on each metric across 15 transits.",
    ("Histograms show CP timing relative to departure (t=0).", True, DARK_GRAY),
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_transit_cp_summary.png", 0.1, 1.4, width=8.0, height=5.8)

add_text(s, 8.3, 1.5, 4.7, 0.4, "Median CUSUM Change-Point",
         size=16, bold=True, color=DARK_BLUE)
add_lines(s, 8.4, 2.0, 4.5, 4.5, [
    ("LHA Gate:      +0.04s   (at departure)", True, GREEN),
    ("LHA Flow:      +0.13s   (at departure)", True, GREEN),
    ("LHA PC1:       +1.33s   (during transit)", False, MED_BLUE),
    ("LHA FR:        +4.34s   (at/after arrival)", False, DARK_GRAY),
    "",
    ("RSP PC1:       +2.54s", False, DARK_GRAY),
    ("RSP FR:        +3.55s", False, DARK_GRAY),
    ("RSP Gate:      +4.27s", False, DARK_GRAY),
    ("RSP Flow:      +4.44s", False, DARK_GRAY),
    "",
    ("LHA dynamics (Gate, Flow) are the first", True, DARK_BLUE),
    ("to shift — essentially coincident with", True, DARK_BLUE),
    ("departure. RSP changes follow 2-4s later.", True, DARK_BLUE),
], size=12, color=DARK_GRAY)


# ─────────────────────────────────────────────────────────
# SLIDE — Destination-dependent divergence
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "Destination-Dependent Divergence: P2->P4 vs P4->P2",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.55, 12.5, 0.8, [
    "Mean trajectory +/- SEM by direction (green = P2->P4, n=9; red = P4->P2, n=6). Bottom panel: Cohen's d.",
    ("LHA PC1 separates by direction before departure. RSP FR consistently higher for P2->P4.", True, DARK_GRAY),
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_transit_divergence_lha.png", 0.1, 1.3, width=6.3, height=6.0)
add_img(s, FIG / "foraging_transit_divergence_rsp.png", 6.5, 1.3, width=6.3, height=6.0)


# ─────────────────────────────────────────────────────────
# SLIDE — Pre-departure shift (negative result)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.1, 12.5, 0.5,
         "Pre-Departure Neural Shift: Early vs Last 2s at Departure Pot",
         size=22, bold=True, color=DARK_BLUE)
add_lines(s, 0.3, 0.55, 12.5, 0.8, [
    "Wilcoxon signed-rank comparing early occupancy vs final 2s before departure. Each line = one transit.",
    ("No significant pre-departure ramp in any metric (all p > 0.12). Decision appears abrupt, not gradual.", True, DARK_GRAY),
], size=11, color=DARK_GRAY)
add_img(s, FIG / "foraging_transit_predeparture.png", 0.2, 1.35, width=12.8, height=5.8)


# ─────────────────────────────────────────────────────────
# SLIDE — Transit dynamics summary
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.15, 12.5, 0.6,
         "Transit Dynamics: Summary", size=24, bold=True, color=DARK_BLUE)
add_lines(s, 0.5, 0.9, 12.0, 6.0, [
    ("1.  Decision is abrupt, not gradual", True, DARK_BLUE),
    "     No pre-departure ramp: early vs last 2s shows no significant shift in any metric",
    "     The neural state does not gradually 'build up' to departure",
    "",
    ("2.  LHA Gate and Flow shift first — right at departure", True, DARK_BLUE),
    "     LHA Gate: median change-point +0.04s from departure (essentially coincident)",
    "     LHA Flow: +0.13s. Other metrics follow 1-4s later during transit/arrival",
    "     Consistent with LHA Gate as the central pot-discrimination signal",
    "",
    ("3.  Neural state encodes current pot identity before departure", True, DARK_BLUE),
    "     LHA PC1 and RSP FR already differ depending on direction (P2->P4 vs P4->P2)",
    "     Cohen's d reaches 0.8-1.5 pre-departure — but confounded with departure pot identity",
    "",
    ("4.  Biggest neural state changes occur during transit (~2s post-departure)", True, DARK_BLUE),
    "     Rate-of-change peaks cluster at +2s from departure across all metrics",
    "     This is when the latent state transitions between pot-specific representations",
], size=12, color=DARK_GRAY)


# ─────────────────────────────────────────────────────────
# SLIDE — Take-home
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.15, 12.5, 0.6,
         "Key Findings", size=28, bold=True, color=DARK_BLUE)

add_lines(s, 0.5, 0.9, 12.0, 6.0, [
    ("1.  LHA encodes pot value via firing rate", True, DARK_BLUE),
    "     S4: LHA FR at Pot-2 drops after food is found elsewhere (p=0.030*)",
    "     Direct pre/post comparison — immune to session drift",
    "",
    ("2.  LHA Gate distinguishes food pot from empty pot", True, DARK_BLUE),
    "     Gate higher at Pot-4 than Pot-2 within the same excursion (p=0.044*, 12/16 transitions)",
    "     S4: 6/6 transitions consistent (p=0.031*). Dynamics 'freeze' at known food location",
    "",
    ("3.  S6 fasted: coordinated LHA-RSP updating at the empty pot", True, DARK_BLUE),
    "     LHA FR/PC1 decrease, RSP PC1/Flow increase over repeated Pot-2 visits",
    "     All 4 signals survive shuffle (p=0.014-0.050*). No session drift. Strongest signal in dataset",
    "",
    ("4.  LHA Flow & Gate at Pot-4 are pot-specific (pre-discovery)", True, DARK_BLUE),
    "     S2: Flow decreases (p=0.020*), Gate increases (p=0.022*) over repeated visits",
    "     Survive shuffle, not explained by drift. Possible encoding of location-specific expectation",
    "",
    ("5.  Decision to leave is abrupt — LHA Gate/Flow shift at departure", True, DARK_BLUE),
    "     No pre-departure ramp. CUSUM change-point: LHA Gate +0.04s, Flow +0.13s from departure",
    "     RSP follows 2-4s later. Biggest state changes occur during transit (~2s post-departure)",
    "",
    ("6.  Fed vs Fasted", True, DARK_BLUE),
    "     Fasted mice find food in <35s with 0 Pot-2 visits — no pre-discovery learning phase",
    "     All pre-discovery signals exclusive to fed sessions; post-discovery signals present in both",
], size=12, color=DARK_GRAY)


# ─────────────────────────────────────────────────────────
# SLIDE 15  — Caveats & next steps
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_text(s, 0.3, 0.15, 12.5, 0.6,
         "Caveats & Next Steps", size=28, bold=True, color=DARK_BLUE)

add_lines(s, 0.5, 0.9, 12.0, 6.0, [
    ("Caveats", True, RED),
    "  N=1 mouse, 2 sessions per condition — findings need replication",
    "  LHA PC1 Pot-2 ramp (previously the 'flagship') does NOT survive shuffle — driven by session-wide drift",
    "  S8 fasted: RSP drift dominates all signals — nothing pot-specific survives",
    "  S2 vs S4 LHA Pot-4 Gate show opposite post-discovery directions — may reflect different strategies",
    "",
    ("Next Steps", True, GREEN),
    "  More mice to confirm LHA Gate as a consistent pot-discriminating signal",
    "  Test whether Gate encodes location identity, value expectation, or certainty",
    "  Decode pot identity from population state at arrival (pre-entry classification)",
    "  Extend to dual-probe recordings with ACA for prefrontal-hypothalamic interactions",
], size=13, color=DARK_GRAY)


# ─────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────
out = Path("data") / "foraging_summary_presentation.pptx"
prs.save(str(out))
print(f"\nSaved: {out}")
print(f"Total slides: {len(prs.slides)}")
